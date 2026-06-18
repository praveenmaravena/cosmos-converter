import os
import io
import shutil
import fitz  # PyMuPDF
from PIL import Image, ImageEnhance, ImageFilter
import pytesseract
from docx import Document as DocxDocument
from pdf2docx import Converter as Pdf2DocxConverter
from pypdf import PdfReader, PdfWriter
import requests
import json
import base64
import zipfile
from app.core.config import settings

class ConversionEngine:
    @staticmethod
    def run_libreoffice_convert(input_path: str, output_dir: str, target_format: str = "pdf") -> bool:
        """Invokes headless LibreOffice as a subprocess to convert office documents."""
        import subprocess
        for cmd in ["soffice", "libreoffice", "C:\\Program Files\\LibreOffice\\program\\soffice.exe"]:
            try:
                subprocess.run(
                    [cmd, "--headless", "--convert-to", target_format, "--outdir", output_dir, input_path],
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL
                )
                return True
            except Exception:
                continue
        return False

    @staticmethod
    def pdf_to_docx(input_path: str, output_path: str):
        cv = Pdf2DocxConverter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()

    @staticmethod
    def docx_to_pdf(input_path: str, output_path: str):
        output_dir = os.path.dirname(output_path)
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        temp_pdf = os.path.join(output_dir, f"{base_name}.pdf")
        
        # Try LibreOffice first
        if ConversionEngine.run_libreoffice_convert(input_path, output_dir, "pdf"):
            if os.path.exists(temp_pdf):
                if temp_pdf != output_path:
                    shutil.move(temp_pdf, output_path)
                return

        # Fallback to python-docx reading and writing using a simplified fitz layout
        doc = DocxDocument(input_path)
        doc_pdf = fitz.open()
        page = doc_pdf.new_page()
        p_y = 50
        for para in doc.paragraphs:
            if p_y > 750:
                page = doc_pdf.new_page()
                p_y = 50
            page.insert_text((50, p_y), para.text, fontsize=11)
            p_y += 20
        doc_pdf.save(output_path)
        doc_pdf.close()

    @staticmethod
    def pdf_to_txt(input_path: str, output_path: str):
        doc = fitz.open(input_path)
        text = ""
        for page in doc:
            text += page.get_text()
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        doc.close()

    @staticmethod
    def txt_to_pdf(input_path: str, output_path: str):
        with open(input_path, "r", encoding="utf-8") as f:
            text = f.read()
        doc = fitz.open()
        page = doc.new_page()
        lines = text.split("\n")
        p_y = 50
        for line in lines:
            if p_y > 750:
                page = doc.new_page()
                p_y = 50
            page.insert_text((50, p_y), line, fontsize=10)
            p_y += 15
        doc.save(output_path)
        doc.close()

    @staticmethod
    def pdf_to_html(input_path: str, output_path: str):
        doc = fitz.open(input_path)
        html = "<html><body>"
        for page in doc:
            html += page.get_text("html")
        html += "</body></html>"
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(html)
        doc.close()

    @staticmethod
    def html_to_pdf(input_path: str, output_path: str):
        with open(input_path, "r", encoding="utf-8") as f:
            html_content = f.read()
        import re
        clean_text = re.sub('<[^<]+?>', '', html_content)
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 50), clean_text[:2000], fontsize=10)
        doc.save(output_path)
        doc.close()

    @staticmethod
    def pdf_to_pptx(input_path: str, output_path: str):
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        doc = fitz.open(input_path)
        blank_slide_layout = prs.slide_layouts[6]
        for page in doc:
            slide = prs.slides.add_slide(blank_slide_layout)
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            text = page.get_text()
            tf.text = text if text.strip() else "[No text on this page]"
        prs.save(output_path)
        doc.close()

    @staticmethod
    def pptx_to_pdf(input_path: str, output_path: str):
        output_dir = os.path.dirname(output_path)
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        temp_pdf = os.path.join(output_dir, f"{base_name}.pdf")
        
        if ConversionEngine.run_libreoffice_convert(input_path, output_dir, "pdf"):
            if os.path.exists(temp_pdf):
                if temp_pdf != output_path:
                    shutil.move(temp_pdf, output_path)
                return

        from pptx import Presentation
        prs = Presentation(input_path)
        doc_pdf = fitz.open()
        for slide in prs.slides:
            page = doc_pdf.new_page()
            p_y = 50
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    if p_y > 750:
                        page = doc_pdf.new_page()
                        p_y = 50
                    page.insert_text((50, p_y), paragraph.text, fontsize=11)
                    p_y += 20
        doc_pdf.save(output_path)
        doc_pdf.close()

    @staticmethod
    def pdf_to_xlsx(input_path: str, output_path: str):
        import openpyxl
        import pdfplumber
        wb = openpyxl.Workbook()
        default_sheet = wb.active
        if default_sheet:
            wb.remove(default_sheet)
        with pdfplumber.open(input_path) as pdf:
            for i, page in enumerate(pdf.pages):
                ws = wb.create_sheet(title=f"Page {i+1}")
                tables = page.extract_tables()
                if tables:
                    row_idx = 1
                    for table in tables:
                        for row in table:
                            for col_idx, val in enumerate(row, start=1):
                                ws.cell(row=row_idx, column=col_idx, value=val)
                            row_idx += 1
                        row_idx += 2
                else:
                    text = page.extract_text() or ""
                    for r_idx, line in enumerate(text.split("\n"), start=1):
                        ws.cell(row=r_idx, column=1, value=line)
        wb.save(output_path)

    @staticmethod
    def xlsx_to_pdf(input_path: str, output_path: str):
        output_dir = os.path.dirname(output_path)
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        temp_pdf = os.path.join(output_dir, f"{base_name}.pdf")
        
        if ConversionEngine.run_libreoffice_convert(input_path, output_dir, "pdf"):
            if os.path.exists(temp_pdf):
                if temp_pdf != output_path:
                    shutil.move(temp_pdf, output_path)
                return

        import openpyxl
        wb = openpyxl.load_workbook(input_path, read_only=True)
        doc_pdf = fitz.open()
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            page = doc_pdf.new_page()
            page.insert_text((50, 30), f"Sheet: {sheet_name}", fontsize=14, color=(0, 0, 0.5))
            p_y = 60
            for row in ws.iter_rows(values_only=True):
                if p_y > 750:
                    page = doc_pdf.new_page()
                    p_y = 50
                row_str = " | ".join([str(val) if val is not None else "" for val in row])
                page.insert_text((50, p_y), row_str[:100], fontsize=9)
                p_y += 15
        doc_pdf.save(output_path)
        doc_pdf.close()

    @staticmethod
    def image_convert(input_path: str, output_path: str, format_name: str):
        img = Image.open(input_path)
        if img.mode in ("RGBA", "P") and format_name.upper() in ("JPG", "JPEG"):
            img = img.convert("RGB")
        img.save(output_path, format=format_name.upper())

    @staticmethod
    def images_to_pdf(input_paths: list, output_path: str):
        images = []
        for path in input_paths:
            img = Image.open(path)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            images.append(img)
        if images:
            images[0].save(output_path, save_all=True, append_images=images[1:])

    @staticmethod
    def pdf_to_single_image(input_path: str, output_path: str, format_name: str = "PNG"):
        doc = fitz.open(input_path)
        if len(doc) > 0:
            page = doc[0]
            pix = page.get_pixmap(dpi=150)
            pix.save(output_path)
        doc.close()

    @staticmethod
    def svg_to_png(input_path: str, output_path: str):
        import subprocess
        for cmd in ["convert", "magick convert", "magick"]:
            try:
                subprocess.run([cmd, input_path, output_path], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                return
            except Exception:
                continue
        try:
            import cairosvg
            cairosvg.svg2png(url=input_path, write_to=output_path)
        except Exception:
            # Fallback placeholder
            img = Image.new("RGB", (400, 400), color=(240, 240, 240))
            img.save(output_path, "PNG")

    @staticmethod
    def png_to_svg(input_path: str, output_path: str):
        with open(input_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
        img = Image.open(input_path)
        width, height = img.size
        ext = os.path.splitext(input_path)[1].replace(".", "").lower()
        if ext == "jpg":
            ext = "jpeg"
        mime = f"image/{ext}"
        svg_content = f'<svg xmlns="http://www.w3.org/2000/svg" xmlns:xlink="http://www.w3.org/1999/xlink" width="{width}" height="{height}" viewBox="0 0 {width} {height}"><image width="{width}" height="{height}" xlink:href="data:{mime};base64,{encoded_string}"/></svg>'
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(svg_content)

    @staticmethod
    def heic_to_image(input_path: str, output_path: str, format_name: str = "JPEG"):
        try:
            from pillow_heif import register_heif_opener
            register_heif_opener()
        except ImportError:
            pass
        img = Image.open(input_path)
        if img.mode in ("RGBA", "P") and format_name.upper() in ("JPG", "JPEG"):
            img = img.convert("RGB")
        img.save(output_path, format=format_name.upper())

    @staticmethod
    def gif_to_mp4(input_path: str, output_path: str):
        import subprocess
        try:
            subprocess.run(
                ["ffmpeg", "-y", "-i", input_path, "-movflags", "faststart", "-pix_fmt", "yuv420p", 
                 "-vf", "scale=trunc(iw/2)*2:trunc(ih/2)*2", output_path],
                check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
            )
        except Exception as e:
            raise RuntimeError(f"FFmpeg GIF to MP4 failed: {e}")

    @staticmethod
    def mp4_to_gif(input_path: str, output_path: str):
        import subprocess
        try:
            subprocess.run(
                ["ffmpeg", "-y", "-i", input_path, "-vf", 
                 "fps=10,scale=480:-1:flags=lanczos,split[s0][s1];[s0]palettegen[p];[s1][p]paletteuse", 
                 "-loop", "0", output_path],
                check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
            )
        except Exception as e:
            raise RuntimeError(f"FFmpeg MP4 to GIF failed: {e}")

    # OCR Operations
    @staticmethod
    def image_to_text(input_path: str) -> str:
        try:
            img = Image.open(input_path)
            return pytesseract.image_to_string(img)
        except Exception as e:
            return f"OCR Error: {str(e)}. Make sure Tesseract-OCR is installed."

    @staticmethod
    def parse_invoice(text: str) -> dict:
        import re
        invoice_data = {
            "invoice_number": None,
            "date": None,
            "total_amount": None,
            "vendor": None
        }
        inv_match = re.search(r'(?:Invoice\s*#?|Inv\s*#?)\s*([A-Za-z0-9-]+)', text, re.IGNORECASE)
        if inv_match:
            invoice_data["invoice_number"] = inv_match.group(1)
        date_match = re.search(r'\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b', text)
        if date_match:
            invoice_data["date"] = date_match.group(0)
        total_match = re.search(r'(?:Total|Amount\s*Due|Balance|Sum):\s*\$?(\d+[\.,]\d{2})', text, re.IGNORECASE)
        if total_match:
            invoice_data["total_amount"] = total_match.group(1)
        return invoice_data

    @staticmethod
    def parse_business_card(text: str) -> dict:
        import re
        card_data = {
            "name": None,
            "email": None,
            "phone": None,
            "website": None
        }
        email_match = re.search(r'[\w\.-]+@[\w\.-]+\.\w+', text)
        if email_match:
            card_data["email"] = email_match.group(0)
        phone_match = re.search(r'\+?\d{1,4}?[-.\s]?\(?\d{1,3}?\)?[-.\s]?\d{1,4}[-.\s]?\d{1,4}[-.\s]?\d{1,9}', text)
        if phone_match:
            card_data["phone"] = phone_match.group(0)
        web_match = re.search(r'(?:www\.|https?://)[\w\.-]+\.\w+', text)
        if web_match:
            card_data["website"] = web_match.group(0)
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        if lines:
            card_data["name"] = lines[0]
        return card_data

    # PDF Advanced Tools
    @staticmethod
    def merge_pdfs(input_paths: list, output_path: str):
        writer = PdfWriter()
        for path in input_paths:
            reader = PdfReader(path)
            for page in reader.pages:
                writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def split_pdf(input_path: str, output_path: str):
        """Splits PDF into separate pages and packs them into a ZIP file."""
        reader = PdfReader(input_path)
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        with zipfile.ZipFile(output_path, 'w') as zipf:
            for i, page in enumerate(reader.pages):
                writer = PdfWriter()
                writer.add_page(page)
                page_pdf_bytes = io.BytesIO()
                writer.write(page_pdf_bytes)
                page_pdf_bytes.seek(0)
                zipf.writestr(f"{base_name}_page_{i+1}.pdf", page_pdf_bytes.read())

    @staticmethod
    def rearrange_pdf(input_path: str, output_path: str, page_order: list):
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for idx in page_order:
            if 0 <= idx < len(reader.pages):
                writer.add_page(reader.pages[idx])
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def rotate_pdf(input_path: str, output_path: str, rotation_angle: int, pages: list = None):
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for idx, page in enumerate(reader.pages):
            if pages is None or idx in pages:
                page.rotate(rotation_angle)
            writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def delete_pdf_pages(input_path: str, output_path: str, pages_to_delete: list):
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for idx, page in enumerate(reader.pages):
            if idx not in pages_to_delete:
                writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def compress_pdf(input_path: str, output_path: str):
        doc = fitz.open(input_path)
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        doc.close()

    @staticmethod
    def protect_pdf(input_path: str, output_path: str, password: str):
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(user_password=password, owner_password=None, use_128bit=True)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def unlock_pdf(input_path: str, output_path: str, password: str):
        reader = PdfReader(input_path)
        if reader.is_encrypted:
            reader.decrypt(password)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)

    @staticmethod
    def add_watermark(input_path: str, output_path: str, watermark_text: str):
        doc = fitz.open(input_path)
        for page in doc:
            rect = page.rect
            width, height = rect.width, rect.height
            page.insert_text(
                fitz.Point(width / 4, height / 2),
                watermark_text,
                fontsize=40,
                color=(0.8, 0.8, 0.8),
                rotate=45,
                fill_opacity=0.3
            )
        doc.save(output_path)
        doc.close()

    @staticmethod
    def add_signature(input_path: str, output_path: str, signature_image_path: str, page_num: int = 0, x: float = 50, y: float = 50, w: float = 150, h: float = 50):
        doc = fitz.open(input_path)
        if 0 <= page_num < len(doc):
            page = doc[page_num]
            rect = fitz.Rect(x, y, x + w, y + h)
            page.insert_image(rect, filename=signature_image_path)
        doc.save(output_path)
        doc.close()

    # Image Processing Toolkit
    @staticmethod
    def resize_image(input_path: str, output_path: str, width: int, height: int):
        img = Image.open(input_path)
        img_resized = img.resize((width, height), Image.Resampling.LANCZOS)
        img_resized.save(output_path)

    @staticmethod
    def crop_image(input_path: str, output_path: str, x: int, y: int, width: int, height: int):
        img = Image.open(input_path)
        img_cropped = img.crop((x, y, x + width, y + height))
        img_cropped.save(output_path)

    @staticmethod
    def rotate_image(input_path: str, output_path: str, angle: float):
        img = Image.open(input_path)
        img_rotated = img.rotate(angle, expand=True)
        img_rotated.save(output_path)

    @staticmethod
    def compress_image(input_path: str, output_path: str, quality: int = 70):
        img = Image.open(input_path)
        fmt = img.format or "JPEG"
        if img.mode in ("RGBA", "P") and fmt.upper() in ("JPG", "JPEG"):
            img = img.convert("RGB")
        img.save(output_path, format=fmt, quality=quality, optimize=True)

    @staticmethod
    def remove_background(input_path: str, output_path: str):
        img = Image.open(input_path).convert("RGBA")
        datas = img.getdata()
        new_data = []
        for item in datas:
            if item[0] > 230 and item[1] > 230 and item[2] > 230:
                new_data.append((255, 255, 255, 0))
            else:
                new_data.append(item)
        img.putdata(new_data)
        img.save(output_path, "PNG")

    @staticmethod
    def ai_enhance_image(input_path: str, output_path: str):
        try:
            import cv2
            import numpy as np
            img_cv = cv2.imread(input_path)
            if img_cv is None:
                raise ValueError()
            lab = cv2.cvtColor(img_cv, cv2.COLOR_BGR2LAB)
            l, a, b = cv2.split(lab)
            clahe = cv2.createCLAHE(clipLimit=3.0, tileGridSize=(8,8))
            cl = clahe.apply(l)
            limg = cv2.merge((cl,a,b))
            enhanced = cv2.cvtColor(limg, cv2.COLOR_LAB2BGR)
            kernel = np.array([[0, -1, 0], [-1, 5,-1], [0, -1, 0]])
            sharpened = cv2.filter2D(enhanced, -1, kernel)
            cv2.imwrite(output_path, sharpened)
        except Exception:
            img = Image.open(input_path)
            enh = ImageEnhance.Contrast(img).enhance(1.2)
            enh = ImageEnhance.Sharpness(enh).enhance(1.3)
            enh.save(output_path)

    @staticmethod
    def adjust_image(input_path: str, output_path: str, brightness: float = 1.0, contrast: float = 1.0, filter_name: str = "none"):
        img = Image.open(input_path)
        if brightness != 1.0:
            enh = ImageEnhance.Brightness(img)
            img = enh.enhance(brightness)
        if contrast != 1.0:
            enh = ImageEnhance.Contrast(img)
            img = enh.enhance(contrast)
        if filter_name == "grayscale":
            img = img.convert("L")
        elif filter_name == "sepia":
            if img.mode != "RGB":
                img = img.convert("RGB")
            width, height = img.size
            pixels = img.load()
            for py in range(height):
                for px in range(width):
                    r, g, b = pixels[px, py]
                    tr = int(0.393 * r + 0.769 * g + 0.189 * b)
                    tg = int(0.349 * r + 0.686 * g + 0.168 * b)
                    tb = int(0.272 * r + 0.534 * g + 0.131 * b)
                    pixels[px, py] = (min(tr, 255), min(tg, 255), min(tb, 255))
        elif filter_name == "blur":
            img = img.filter(ImageFilter.BLUR)
        elif filter_name == "contour":
            img = img.filter(ImageFilter.CONTOUR)
        elif filter_name == "sharpen":
            img = img.filter(ImageFilter.SHARPEN)
        img.save(output_path)

